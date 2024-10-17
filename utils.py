from pptx import Presentation
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
import os
import fitz
from pptxtopdf import convert
import win32com.client as win32
import json

# I-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-I

config_path = 'config.json'  # Arquivo de configuração, para que exista um meio do usuario mudar algumas coisas


# Carregamento e leitura desse arquivo
def load_config(config_p):
    if os.path.exists(config_p):
        with open(config_p, 'r') as file:
            return json.load(file)
    else:
        return {  # Se não existir dados, ele usa esses aqui.
            "horario": "15:00",
            "email": "email_default@default.com",
            "planilha": "planilha_cliente.xlsx",
            "pptx0_file": "Informativo - Email Especialistas.pptx"
        }


config = load_config(config_path)  # atribui a variavel o dicionario do arquivo json
# I-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-I

x = 0  # Variavel global que ajuda a função "make_slide()" a manter os indices das listas

wb = load_workbook(config['planilha'])  # Carrega o arquivo .xlsx determinado na config.json
ws = wb.active  # deixa o "workbook", a planilha, ativa.

lista = []  # cria uma lista que irá conter todos os dados da planilha
lista_copy = lista  # Faz uma cópia da lista pois a mesma irá ser desfeita, pelo pop()
path = config['pptx0_file']  # Armazena o caminho do template do pptx
diretorio = os.path.join(os.getcwd(), 'Slides')  # Armazena na variavel o caminho que será armazenado os pdfs e afins

# Carrega o outlook
outlook = win32.Dispatch("Outlook.Application")
namespace = outlook.GetNamespace("MAPI")

# Carrega os slides pptx
ppt = Presentation(path)
slide = ppt.slides[0]  # Pega o unico slide

# Essas duas listas são feitas para separar o nome e o email de cada funcionario
# Seram usada na função send_email()

email_list = []
name_list = []

# Preenche a lista com todas as informaçÕes da planilha
# separando cada linha da planilha, em uma lista, fazendo uma lista de listas

for row in ws.iter_rows(values_only=True):
    if any(cell is not None for cell in row):
        lista.append(row)

# Faz a mesma coisa que a anterior, mas preenche as listas:
# email_funcionario e name_funcionario

for indice, tupla in enumerate(lista):
    email_funcionario = lista[indice][-1].split(":")[1].replace(" ", "")
    name_funcionario = lista_copy[indice][-1].split(":")[0]
    if " " in name_funcionario:
        name_funcionario = name_funcionario.split(" ")[0]
    email_list.append(email_funcionario)
    name_list.append(name_funcionario)

# I-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=Cemitério-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-I

"""
Abaixo jás uma função que pediram para não implementar: 
deixarei aqui via as duvidas


# Cheva se "Default" se encontra no primeiro bloco do Excel
def check0():
    if 'Defaut' in ws['A1'].value:
        # se sim, retorne True
        return True
    # Se não, retorne False
    return False

# Faz um template inicial para a planilha, deixando-a inteira em "Default-Mode"
def default_p():
    ws.delete_rows(1, ws.max_row)  # Deleta a planilha inteira 
    dados_default = ['Empresa_Default', '11.111.111/0000-00', '(11) 11111-1111', 'contato@default.com', 'nome_default',
                     'Funcionario Default: funcionario_default@gmail.com']
    for coluna, valor in enumerate(dados_default, 1):
        ws.cell(1, coluna, valor)  # Insere os dados "Default"
    wb.save(config["planilha"])  # Salva a planilha

"""


# I-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-I

#

class Cliente:
    """
    Classe cliente:
    Armazena alguns dados importantes para serem chamados mais tarde
    """

    def __init__(self, empresa, cnpj, celular, email, nome_s):
        self.empresa = empresa
        self.cnpj = cnpj
        self.celular = celular
        self.email = email
        self.nome_s = nome_s


def list_s():
    """
    Retorna as caixas de texto que existem nos arquivos pptx
    :return: list de caixas_de_texto de um pptx
    """
    text_boxes = []
    for shape in slide.shapes:  # Para cada layer no slide, verifica se existe caixas de texto
        if shape.has_text_frame and shape.text:  # Se a camada tem caixa de texto E tem texto nela, adiciona na lista
            text_boxes.append(shape.text)
    return text_boxes


def change_text(txt_id: int, new_txt: str):
    """
    Funciona em dois turnos:
    No primeiro; ele armazena o padrão do texto anterior.
    No segundo; Ele aplica o texto novo, pedindo o ID do texto e o novo texto em si.
    :param txt_id: int.
    :param new_txt: str.
    :return: nothing
    """
    count = 0
    for shape in slide.shapes:  # Para cada layer no slide:
        if shape.has_text_frame and shape.text:  # verifica se existe caixas de texto nessa layer

            # Acrescenta um numero ao contador, que verifica se essa caixa de texto é a mesma fornecida no ID
            count += 1
            if count == txt_id:  # se for o texto escolhido pela função,
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
                # Armazena a formatação da primeira run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                # Deleta todo e qualquer texto anterior e aplica o texto novo:

                text_frame.clear()  # Deleta todo o texto e sua formatação
                new_run = text_frame.paragraphs[0].add_run()  # Cria o novo texto no primeiro paragrafo
                new_run.text = new_txt
                # Aplica a nova run
                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                if count == 5:
                    new_run.font.color.rgb = RGBColor(255, 255, 255)
                else:
                    new_run.font.color.rgb = RGBColor(0, 0, 0)  # Aplica a cor dos textos. No caso, branco.
                return


async def make_slide():
    global x
    nome_funcionario = lista_copy[0][5].split(":")[0]
    if " " in nome_funcionario:
        nome_funcionario = nome_funcionario.split(" ")[0]

    cliente = Cliente(lista_copy[0][0], lista_copy[0][1], lista_copy[0][2], lista_copy[0][3], lista_copy[0][4])

    change_text(5, nome_funcionario + ',')
    change_text(8, cliente.empresa)
    change_text(7, cliente.cnpj)
    change_text(9, cliente.celular)
    change_text(10, cliente.email)
    change_text(16, cliente.nome_s)

    file = os.path.join(diretorio, f'Email_{str(x + 1)}_Especialistas.pptx')
    ppt.save(file)
    x += 1
    print(f'Slides: {x} PRONTO!')

    lista_copy.pop(0)
    return file


def clear_files(file: str):
    """
    Deleta um arquivo
    :param file: str. Path do arquivo escolhido
    :return: nothing
    """
    try:
        os.remove(file)
        print(f"Deleted: {file}")
    except Exception as error:
        print(f"Failed to delete {file}: {error}")  # Se não deletado, retona um erro e o arquivo


async def convert_to_pdf(file: str):
    """
    Converte todos os arquivos pptx, e retorna o True para manter a fila rodando.
    :param file: str. Path do arquivo a ser convertido
    :return: bool.
    """
    try:
        convert(diretorio, diretorio)  # Converte todos os arquivos da pasta Slides, e os converte lá mesmo.
        print("Conversion done!")
        clear_files(file)  # Limpa o arquivo convertido
        return True
    except Exception as e:
        print(f"(PPTX) Conversion failed: {e}")


async def convert_to_img(file: str):
    """
    Converte o arquivo passado .pdf, e retorna o arquivo convertido para manter a fila rodando.
    :param file: str. Path do arquivo a converter de PDF para JPG
    :return: str. Path do arquivo já convertido
    """
    try:
        if ".pdf" in file:
            # Abre o pdf, carrega a unica pagina, converte o arquivo, armazena a nova extensão em file_v2,
            # salva, fecha, e exclui o arquivo anterior .pdf.
            pdf = fitz.open(file)
            page = pdf.load_page(0)
            pix = page.get_pixmap()
            file_v2 = file.replace(".pdf", ".jpg")
            pix.save(file_v2)
            pdf.close()
            clear_files(file)
            return file_v2
    except Exception as e:
        print(f"(PDF) Conversion failed: {e}")


def check_conta(mail):
    """
    Verifica se a conta (colocada no config.json) existe, se existir, retorna True e muda para ela.
    :param mail: class.
    :return: bool. Retorna True para sinalizar que a conta foi mudada.
    """
    for myEmailAddress in outlook.Session.Accounts:
        if config['email'] in str(myEmailAddress):
            mail._oleobj_.Invoke(*(64209, 0, 8, 0, myEmailAddress))
            return True


async def send_email(file):
    """
    Função de enviar email, recebe um parametro que é o jpg que vai ser utilizado para anexar ao email.
    :param file: str. Path do arquivo JPG.
    :return: nothing
    """
    nome_funcionario = name_list.pop(0)  # Armazena o nome do funcionario ao mesmo tempo que o tira da lista/fila.
    mail = outlook.CreateItem(0)  # Cria um objeto de email
    if check_conta(mail):  # Checa se existe a conta colocada em config.json
        print("email aceito")
        mail.Subject = f"{nome_funcionario}, Veja seus clientes!!"  # Assunto do email
        mail.HTMLBody = f"""<html>
                    <body>
                        <h2 style="border: 2px solid black; padding: 10px; color: white; text-align: center;">
                            {nome_funcionario}, Abra para ver seus clientes:
                        </h2>
                        <img src="cid:image1" width="500" style="display: block; margin: 0 auto;">
                    </body>
                    </html>
                    """

        attachment = mail.Attachments.Add(file)  # Coloca o anexo, no caso a imagem enviada.
        attachment.PropertyAccessor.SetProperty("http://schemas.microsoft.com/mapi/proptag/0x3712001F",
                                                "image1")  # Envia o anexo no formado desejado. (Inteiro/Descompactado)
        attachment_v2 = os.path.join(os.getcwd(), config["anexo_2"])
        mail.Attachments.Add(attachment_v2)
        mail.To = email_list[0]  # Envia para o email que está na lista
        try:
            mail.Send()  # Envia
            print('Enviado para:', email_list[0])
            email_list.pop(0)  # Remove o email já enviado para enviar o proximo
        except Exception as error:
            print('Não enviado: ', error)
        clear_files(file)  # Deleta o arquivo do email já enviado

    #   default_p()  Função não implementada que deixa a planilha principal em Default
